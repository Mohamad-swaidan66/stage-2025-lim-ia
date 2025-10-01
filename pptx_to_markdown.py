from pptx import Presentation
from pathlib import Path


def extract_text_from_shape(shape):
    """
    Extrait récursivement le texte d’un shape (y compris groupes et tableaux).
    """
    text = ""

    if hasattr(shape, "text"):
        text += shape.text.strip() + "\n"

    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            text += extract_text_from_shape(sub_shape)

    if shape.has_table:
        for row in shape.table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text += row_text + "\n"

    return text


def pptx_to_markdown(pptx_path, output_file=None):
    """
    Convertit un fichier PowerPoint (.pptx) en Markdown structuré.
    """
    prs = Presentation(pptx_path)
    md = f"# Contenu du fichier : {pptx_path.name}\n\n"

    for i, slide in enumerate(prs.slides, start=1):
        md += f"## Slide {i}\n\n"
        for shape in slide.shapes:
            extracted = extract_text_from_shape(shape)
            if extracted.strip():
                # Nettoyage des lignes vides superflues
                lines = [line.strip() for line in extracted.splitlines() if line.strip()]
                md += "\n".join(lines) + "\n\n"

    if output_file:
        output_file = Path(output_file)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(md)
        print(f"✅ Fichier Markdown sauvegardé : {output_file}")

    return md


if __name__ == "__main__":
    # Chemin d'entrée du fichier PowerPoint
    pptx_path = Path("var/www/RAG/code/utils/document_converters/docling_pdf_to_markdown.py")  # ← ton fichier ici

    # Chemin de sortie pour le fichier Markdown
    output_md_path = Path("/var/www/RAG/Data_parse/powerpoint.md")

    # Conversion
    markdown_result = pptx_to_markdown(pptx_path, output_file=output_md_path)

    # Affichage dans le terminal (aperçu)
    print("\n--- Aperçu Markdown ---\n")
    print(markdown_result[:1000000]) 